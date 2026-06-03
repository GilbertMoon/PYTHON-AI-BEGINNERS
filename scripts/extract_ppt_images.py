"""
4개 강의 PPT 파일에서 슬라이드에 표시된 이미지를 추출하여 images/ 폴더에 저장합니다.

사용 방법
1. 레포지토리 루트 또는 PPT/ 폴더에 PPT 파일 4개를 둡니다.
   파일명은 PART1, PART2, PART3, PART4로 시작하면 자동으로 찾습니다.
2. 필요한 패키지를 설치합니다.
   python -m pip install --upgrade pymupdf python-pptx pillow
3. 실행합니다.
   python scripts/extract_ppt_images.py

참고
- 일반 PNG/JPG 이미지는 PPT 내부 crop 값을 반영하여 저장합니다.
- WMF처럼 Pillow가 직접 읽지 못하는 이미지는 LibreOffice로 슬라이드를 PDF 렌더링한 뒤,
  슬라이드에 표시된 위치 기준으로 잘라 저장합니다.
- WMF 처리까지 하려면 PC에 LibreOffice가 설치되어 있고, soffice/libreoffice 명령이 PATH에 있어야 합니다.
"""

from __future__ import annotations

import csv
import io
import shutil
import subprocess
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT_DIR = Path(__file__).resolve().parents[1]
PPT_DIR = ROOT_DIR / "PPT"
IMAGE_DIR = ROOT_DIR / "images"
RENDER_DIR = ROOT_DIR / "_rendered_pdfs"

PPT_PATTERNS = [
    ("chapter1", "PART1*.pptx"),
    ("chapter2", "PART2*.pptx"),
    ("chapter3", "PART3*.pptx"),
    ("chapter4", "PART4*.pptx"),
]


def find_ppt_file(pattern: str) -> Path | None:
    """ROOT_DIR와 PPT_DIR에서 pattern에 맞는 PPT 파일을 찾습니다."""
    candidates: list[Path] = []

    for base_dir in (PPT_DIR, ROOT_DIR):
        if base_dir.exists():
            candidates.extend(sorted(base_dir.glob(pattern)))

    # 임시 파일(~$...) 제외
    candidates = [path for path in candidates if not path.name.startswith("~$")]

    if not candidates:
        return None

    return candidates[0]


def clamp(value: float, min_value: float, max_value: float) -> float:
    return max(min_value, min(max_value, value))


def convert_ppt_to_pdf(ppt_path: Path) -> Path:
    RENDER_DIR.mkdir(exist_ok=True)
    pdf_path = RENDER_DIR / f"{ppt_path.stem}.pdf"
    if pdf_path.exists() and pdf_path.stat().st_size > 0:
        return pdf_path

    command = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(RENDER_DIR),
        str(ppt_path),
    ]

    try:
        subprocess.run(command, check=True)
    except FileNotFoundError:
        command[0] = "soffice"
        subprocess.run(command, check=True)

    return pdf_path


def save_blob_cropped_image(shape, output_path: Path) -> tuple[int, int, str]:
    image = Image.open(io.BytesIO(shape.image.blob)).convert("RGBA")
    width, height = image.size

    left = int(round(clamp(shape.crop_left or 0, 0, 1) * width))
    top = int(round(clamp(shape.crop_top or 0, 0, 1) * height))
    right = int(round(width - clamp(shape.crop_right or 0, 0, 1) * width))
    bottom = int(round(height - clamp(shape.crop_bottom or 0, 0, 1) * height))

    cropped = image if right <= left or bottom <= top else image.crop((left, top, right, bottom))
    cropped.save(output_path, "PNG", optimize=True)
    return cropped.width, cropped.height, "blob-crop"


def save_rendered_cropped_image(shape, output_path: Path, prs, pdf_doc, slide_no: int) -> tuple[int, int, str]:
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    page = pdf_doc[slide_no - 1]
    pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
    slide_image = Image.frombytes("RGB", [pixmap.width, pixmap.height], pixmap.samples)

    left = int(round(shape.left / slide_width * pixmap.width))
    top = int(round(shape.top / slide_height * pixmap.height))
    right = int(round((shape.left + shape.width) / slide_width * pixmap.width))
    bottom = int(round((shape.top + shape.height) / slide_height * pixmap.height))

    left = max(0, min(pixmap.width, left))
    right = max(0, min(pixmap.width, right))
    top = max(0, min(pixmap.height, top))
    bottom = max(0, min(pixmap.height, bottom))

    if right <= left or bottom <= top:
        raise ValueError("이미지 영역을 계산할 수 없습니다.")

    cropped = slide_image.crop((left, top, right, bottom))
    cropped.save(output_path, "PNG", optimize=True)
    return cropped.width, cropped.height, "slide-render-crop"


def main() -> None:
    if IMAGE_DIR.exists():
        shutil.rmtree(IMAGE_DIR)
    IMAGE_DIR.mkdir(parents=True)

    manifest = []

    for chapter_name, pattern in PPT_PATTERNS:
        ppt_path = find_ppt_file(pattern)

        if ppt_path is None:
            print(f"[SKIP] PPT 파일 없음: {pattern}")
            continue

        print(f"[START] {chapter_name}: {ppt_path.relative_to(ROOT_DIR)}")
        chapter_dir = IMAGE_DIR / chapter_name
        chapter_dir.mkdir(parents=True, exist_ok=True)

        prs = Presentation(str(ppt_path))
        pdf_doc = None

        for slide_no, slide in enumerate(prs.slides, start=1):
            image_no = 0
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue

                image_no += 1
                file_name = f"{chapter_name}_slide{slide_no:03d}_img{image_no:02d}.png"
                output_path = chapter_dir / file_name

                try:
                    width, height, source = save_blob_cropped_image(shape, output_path)
                except Exception:
                    if pdf_doc is None:
                        pdf_path = convert_ppt_to_pdf(ppt_path)
                        pdf_doc = fitz.open(str(pdf_path))
                    width, height, source = save_rendered_cropped_image(shape, output_path, prs, pdf_doc, slide_no)

                manifest.append(
                    {
                        "chapter": chapter_name,
                        "slide_no": slide_no,
                        "image_no": image_no,
                        "path": str(output_path.relative_to(IMAGE_DIR)).replace("\\", "/"),
                        "width": width,
                        "height": height,
                        "source": source,
                    }
                )

        if pdf_doc is not None:
            pdf_doc.close()

        print(f"[DONE] {chapter_name}")

    manifest_path = IMAGE_DIR / "manifest.csv"
    with manifest_path.open("w", encoding="utf-8", newline="") as file:
        fieldnames = ["chapter", "slide_no", "image_no", "path", "width", "height", "source"]
        writer = csv.DictWriter(file, fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(manifest)

    print(f"[COMPLETE] 총 {len(manifest)}개 이미지 추출 완료")
    print(f"[MANIFEST] {manifest_path}")


if __name__ == "__main__":
    main()
